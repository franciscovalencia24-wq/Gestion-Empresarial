import os
import streamlit as st
import pandas as pd
import datetime
from src.database.connection import SessionLocal
from src.database.models import MarketVision
from src.agents.macro_analyst import MacroAnalystAgent
from src.osint.institutional_scraper import InstitutionalScraper

def render_macro_chat_ui():
    st.markdown("""
        <div style='background: linear-gradient(135deg, #0f172a 0%, #1e3a5f 100%);
                    padding: 28px; border-radius: 15px; margin-bottom: 25px; color: white;'>
            <h1 style='color: white; margin: 0; font-size: 2.2em; font-weight: 900;'>
                🌍 Analista Macro (Dashboard Institucional)
            </h1>
            <p style='color: #94a3b8; margin: 8px 0 0 0; font-size: 1.1em;'>
                Cuadrícula temporal de visiones de mercado y Consenso IA.
            </p>
        </div>
    """, unsafe_allow_html=True)

    # Opciones de Período Histórico
    from dateutil.relativedelta import relativedelta
    today = datetime.date.today()
    periodos_historicos = [(today.replace(day=1) - relativedelta(months=i)).strftime('%Y-%m') for i in range(12)]
    
    colA, colB, colC = st.columns([1,1,2])
    with colA:
        periodo_actual = st.selectbox("📅 Seleccionar Período", periodos_historicos, index=0)
    with colB:
        st.write("")
        st.write("")
        if st.button("🤖 Generar Consenso IA para este Mes", type="primary", use_container_width=True):
            st.session_state.trigger_consensus = True
            
    st.markdown("---")
    
    tab1, tab2, tab3, tab4 = st.tabs(["📊 Cuadrícula de Instituciones", "📥 Ingesta (PDF/Web)", "💬 Chat Omnicanal", "🔬 Inteligencia & Playbooks"])

    with tab1:
        st.subheader(f"Visiones de Mercado - {periodo_actual}")
        
        db = SessionLocal()
        visiones = db.query(MarketVision).filter_by(periodo=periodo_actual).all()
        db.close()
        
        if not visiones:
            st.info(f"No hay visiones de mercado registradas para {periodo_actual}. Ve a la pestaña 'Ingesta' para cargar información.")
        else:
            # Mostrar como cards/grilla
            cols = st.columns(3)
            for idx, v in enumerate(visiones):
                with cols[idx % 3]:
                    st.markdown(f"""
                    <div style="background-color: #1e1e1e; padding: 15px; border-radius: 8px; border-top: 4px solid #3b82f6; margin-bottom: 15px;">
                        <h4 style="margin-top:0; color: #3b82f6;">🏛️ {v.institucion}</h4>
                        <p style="font-size: 0.8rem; color: #9ca3af;">Fuente: {v.fuente} | {v.fecha_ingesta.strftime('%d/%m/%Y')}</p>
                        <p style="font-size: 0.95rem; color: #e5e7eb;">{v.resumen_corto}</p>
                    </div>
                    """, unsafe_allow_html=True)
                    with st.expander("Ver Visión Extendida"):
                        st.write(v.resumen_extendido)

        # Cargar Consenso guardado desde DB o procesar nuevo si fue disparado
        if getattr(st.session_state, "trigger_consensus", False):
            st.session_state.trigger_consensus = False
            with st.spinner(f"Altus AI Macro está leyendo las visiones de {periodo_actual} y debatiendo el consenso..."):
                agent = MacroAnalystAgent()
                consenso_text = agent.generate_monthly_consolidation(periodo_actual)
                st.session_state[f"last_consensus_{periodo_actual}"] = consenso_text

        if not st.session_state.get(f"last_consensus_{periodo_actual}"):
            db_c = SessionLocal()
            c_saved = db_c.query(MarketVision).filter_by(institucion="Consenso IA", periodo=periodo_actual).first()
            if c_saved and c_saved.resumen_extendido:
                st.session_state[f"last_consensus_{periodo_actual}"] = c_saved.resumen_extendido
            db_c.close()

        current_consensus = st.session_state.get(f"last_consensus_{periodo_actual}")
        if current_consensus:
            st.markdown("---")
            st.markdown(f"## 🧠 Consenso Definitivo del Mes ({periodo_actual})")
            st.success(current_consensus)

            # Sección de Descarga Exportable Institucional
            st.markdown("### 📥 Descargar Informe Institucional Exportable")
            col_pdf, col_docx = st.columns(2)
            
            import tempfile
            temp_dir = tempfile.gettempdir()

            with col_pdf:
                try:
                    import importlib
                    import src.utils.pdf_generator_macro
                    importlib.reload(src.utils.pdf_generator_macro)
                    from src.utils.pdf_generator_macro import generate_macro_pdf

                    pdf_path = os.path.join(temp_dir, f"Consenso_Institucional_{periodo_actual}.pdf")
                    generate_macro_pdf(
                        cliente_nombre="Comité de Inversiones / Clientes FV",
                        contenido_markdown=current_consensus,
                        output_path=pdf_path
                    )
                    with open(pdf_path, "rb") as f_pdf:
                        st.download_button(
                            label="📄 Descargar Informe Completo en PDF",
                            data=f_pdf.read(),
                            file_name=f"Informe_Consenso_Macro_{periodo_actual}.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )
                except Exception as e_pdf:
                    st.error(f"Error generando PDF: {e_pdf}")

            with col_docx:
                try:
                    import importlib
                    import src.utils.docx_generator_macro
                    importlib.reload(src.utils.docx_generator_macro)
                    from src.utils.docx_generator_macro import generate_macro_docx

                    docx_path = os.path.join(temp_dir, f"Consenso_Institucional_{periodo_actual}.docx")
                    generate_macro_docx(
                        cliente_nombre="Comité de Inversiones / Clientes FV",
                        contenido_markdown=current_consensus,
                        output_path=docx_path
                    )
                    with open(docx_path, "rb") as f_docx:
                        st.download_button(
                            label="📝 Descargar Informe Completo en Word (.docx)",
                            data=f_docx.read(),
                            file_name=f"Informe_Consenso_Macro_{periodo_actual}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True
                        )
                except Exception as e_docx:
                    st.error(f"Error generando Word: {e_docx}")

            st.markdown("---")
            st.markdown("### 🌐 Publicación Web Automática")
            if st.button("🚀 Publicar Consenso en fv-inversiones.com", use_container_width=True):
                with st.spinner("Conectando con Vercel y actualizando repositorio..."):
                    import src.utils.web_publisher
                    import importlib
                    importlib.reload(src.utils.web_publisher)
                    from src.utils.web_publisher import publish_consensus_to_vercel
                    success, msg = publish_consensus_to_vercel(current_consensus, periodo_actual)
                    if success:
                        st.success(msg)
                    else:
                        st.error(f"Error publicando: {msg}")

    with tab2:
        c1, c2 = st.columns(2)
        with c1:
            st.subheader("🌐 Ingesta Web (Robots OSINT)")
            st.write("Dispara los robots para extraer la visión de las páginas públicas (SURA, Banchile, Santander, etc.)")
            if st.button("Ejecutar Scrapers Institucionales", use_container_width=True):
                with st.spinner(f"Desplegando robots web OSINT para {periodo_actual}..."):
                    import importlib
                    import src.osint.institutional_scraper
                    importlib.reload(src.osint.institutional_scraper)
                    from src.osint.institutional_scraper import InstitutionalScraper

                    scraper = InstitutionalScraper()
                    try:
                        res = scraper.run_all_scrapers(periodo_actual)
                    except TypeError:
                        try:
                            res = scraper.run_all_scrapers(target_period=periodo_actual)
                        except Exception:
                            res = scraper.run_all_scrapers()

                    st.success(f"✅ Scrapers ejecutados. Las visiones de {periodo_actual} han sido guardadas y consolidadas.")
                    if isinstance(res, dict) and res.get("alerts"):
                        st.warning("Se detectaron alertas en algunas páginas: " + str(res["alerts"]))
                    import time
                    time.sleep(2.5)
                    st.rerun()

        with c2:
            st.subheader("📄 Ingesta de Documentos o URLs (JP Morgan)")
            st.write("Sube PDFs/PPTs largos o pega una URL interactiva para que el Robot le tome foto a los gráficos.")
            inst_name = st.text_input("Nombre de la Institución:")
            mes_pdf = st.selectbox("¿A qué período corresponde?", periodos_historicos, index=0)
            
            tipo_doc = st.radio("Tipo de Análisis", ["Estándar (Texto/Archivo)", "Multimodal Avanzado (JP Morgan / Archivo)", "Multimodal desde URL (Pantallazo Autónomo)"], help="El Multimodal Avanzado le permite a la IA 'ver' los gráficos. La opción URL abrirá un navegador invisible para capturar la web.")
            
            doc_file = None
            url_input = ""
            if tipo_doc == "Multimodal desde URL (Pantallazo Autónomo)":
                url_input = st.text_input("Pega la URL del reporte interactivo (ej. JP Morgan Guide to the Markets):")
            else:
                doc_file = st.file_uploader("Seleccionar Archivo", type=['pdf', 'pptx'])
            
            if st.button("Procesar Ingesta con IA", use_container_width=True):
                if not inst_name:
                    st.warning("Debes indicar el nombre de la institución.")
                elif tipo_doc != "Multimodal desde URL (Pantallazo Autónomo)" and doc_file is None:
                    st.warning("Debes adjuntar el archivo.")
                elif tipo_doc == "Multimodal desde URL (Pantallazo Autónomo)" and not url_input:
                    st.warning("Debes ingresar una URL válida.")
                else:
                    import importlib
                    import src.agents.macro_analyst
                    importlib.reload(src.agents.macro_analyst)
                    agent = src.agents.macro_analyst.MacroAnalystAgent()
                    
                    if tipo_doc == "Estándar (Texto/Archivo)":
                        text_content = ""
                        with st.spinner("Extrayendo texto del documento..."):
                            if doc_file.name.endswith('.pdf'):
                                import PyPDF2
                                pdf_reader = PyPDF2.PdfReader(doc_file)
                                for page in pdf_reader.pages:
                                    text_content += page.extract_text() + "\n"
                            elif doc_file.name.endswith('.pptx'):
                                from pptx import Presentation
                                prs = Presentation(doc_file)
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            text_content += shape.text + "\n"
                                            
                        with st.spinner(f"Resumiendo visión de {inst_name} para {mes_pdf}..."):
                            res = agent.ingest_document_vision(inst_name, mes_pdf, text_content=text_content)
                            if "✅" in res:
                                st.success(res)
                                import time; time.sleep(2.5)
                                st.rerun()
                            else:
                                st.error(res)
                                
                    elif tipo_doc == "Multimodal Avanzado (JP Morgan / Archivo)":
                        with st.spinner(f"Subiendo documento a la visión multimodal de Gemini para {inst_name}..."):
                            temp_path = os.path.join(os.getcwd(), "temp_" + doc_file.name)
                            with open(temp_path, "wb") as f:
                                f.write(doc_file.getbuffer())
                                
                            res = agent.ingest_document_vision(inst_name, mes_pdf, file_path=temp_path, is_multimodal=True)
                            
                            try:
                                os.remove(temp_path)
                            except:
                                pass
                                
                            if "✅" in res:
                                st.success(res)
                                import time; time.sleep(2.5)
                                st.rerun()
                            else:
                                st.error(res)
                                
                    elif tipo_doc == "Multimodal desde URL (Pantallazo Autónomo)":
                        from src.osint.screenshot_engine import capture_full_page_screenshot
                        with st.spinner(f"El Robot Invisible está navegando a la URL y tomando pantallazos... (Esto puede tomar hasta 1 minuto)"):
                            try:
                                temp_path = capture_full_page_screenshot(url_input)
                                
                                with st.spinner("Enviando pantallazo a Gemini para análisis visual de gráficos..."):
                                    res = agent.ingest_document_vision(inst_name, mes_pdf, file_path=temp_path, is_multimodal=True)
                                    
                                    try:
                                        os.remove(temp_path)
                                    except:
                                        pass
                                        
                                    if "✅" in res:
                                        st.success(res)
                                        import time; time.sleep(2.5)
                                        st.rerun()
                                    else:
                                        st.error(res)
                            except Exception as e:
                                st.error(f"Error al capturar la pantalla: {str(e)}")

    with tab3:
        st.write("Escribe tu consulta y el Analista Macro te responderá de forma fundamentada.")
        
        with st.expander("📎 Adjuntar Contexto Externo (YouTube / Noticia)"):
            url = st.text_input("Pega el enlace aquí para que la IA lo use de contexto:")
            contexto_extra = f"[Contenido de {url}]" if url else ""
            
        query = st.text_area("¿Qué deseas preguntarle a Altus AI Macro?", height=150, placeholder="Ej: ¿Qué opina la industria chilena sobre el recorte de tasas?")
        
        if st.button("Enviar Consulta Macro", type="primary"):
            if not query:
                st.warning("Debes escribir una consulta.")
            else:
                with st.spinner("Analizando tu consulta..."):
                    agent = MacroAnalystAgent()
                    current_rut = st.session_state.get("current_client_rut")
                    respuesta = agent.answer_macro_query(query, contexto_extra, client_rut=current_rut)
                    st.session_state["macro_last_response"] = respuesta
                    st.session_state["macro_last_query"] = query
                    
        if st.session_state.get("macro_last_response"):
            st.markdown("### 🤖 Respuesta (Editable):")
            edited_response = st.text_area(
                "Edita la respuesta antes de generar el reporte si lo deseas:", 
                value=st.session_state["macro_last_response"], 
                height=400,
                key="edited_macro_response_area"
            )
            
            c_save, c_dl = st.columns(2)
            with c_save:
                if st.button("💾 Guardar en Historial del Cliente", use_container_width=True):
                    current_rut = st.session_state.get("current_client_rut")
                    last_query = st.session_state.get("macro_last_query", "Consulta no registrada")
                    if current_rut:
                        db = SessionLocal()
                        from src.database.models import Prospect, ClientMacroHistory
                        prospect = db.query(Prospect).filter_by(rut=current_rut).first()
                        if prospect:
                            nuevo_historial = ClientMacroHistory(
                                prospect_id=prospect.id,
                                pregunta=last_query,
                                respuesta_final=edited_response
                            )
                            db.add(nuevo_historial)
                            db.commit()
                            st.success("✅ ¡Guardado en el historial exitosamente!")
                        db.close()
                    else:
                        st.warning("No hay un cliente activo seleccionado para guardar el historial.")

            with c_dl:
                # Generar reporte PDF (Un solo botón)
                import src.utils.pdf_generator_macro
                import importlib
                importlib.reload(src.utils.pdf_generator_macro)
                from src.utils.pdf_generator_macro import generate_macro_pdf
                import tempfile
                
                client_name = st.session_state.get("current_client_name", "Cliente")
                temp_pdf_path = os.path.join(tempfile.gettempdir(), f"Macro_Report_FV.pdf")
                
                try:
                    generate_macro_pdf(client_name, edited_response, temp_pdf_path)
                    with open(temp_pdf_path, "rb") as f:
                        pdf_bytes = f.read()
                    
                    st.download_button(
                        label="📥 Generar y Descargar Reporte Institucional (PDF)",
                        data=pdf_bytes,
                        file_name=f"Reporte_Macro_{client_name.replace(' ', '_')}.pdf",
                        mime="application/pdf",
                        type="primary"
                    )
                except Exception as e:
                    st.error(f"Error generando PDF: {e}")

        st.markdown("---")
        st.subheader("📚 Historial de Consultas del Cliente")
        current_rut = st.session_state.get("current_client_rut")
        if current_rut:
            db = SessionLocal()
            from src.database.models import Prospect, ClientMacroHistory
            try:
                prospect = db.query(Prospect).filter_by(rut=st.session_state["current_client_rut"]).first()
                if prospect:
                    # 🔗 Buscar historial propio y de entidades relacionadas
                    related = prospect.get_related_prospects(db)
                    related_ids = [r.id for r in related]
                    all_ids = [prospect.id] + related_ids
                    
                    historial = db.query(ClientMacroHistory).filter(
                        ClientMacroHistory.prospect_id.in_(all_ids)
                    ).order_by(ClientMacroHistory.fecha.desc()).all()
                    
                    # Mostrar historial Macro
                    if not historial:
                        st.info("No hay historial de chat guardado para este cliente ni sus entidades relacionadas.")
                    else:
                        for h in historial:
                            with st.expander(f"Consulta: {h.fecha.strftime('%Y-%m-%d %H:%M')}"):
                                # Indicador de origen si viene de otra entidad
                                if h.prospect_id != prospect.id:
                                    owner = db.query(Prospect).filter_by(id=h.prospect_id).first()
                                    st.caption(f"🔄 **Respondido originalmente para:** {owner.nombre if owner else 'Desconocido'} ({owner.rut if owner else ''})")
                                
                                st.markdown("**Pregunta:**")
                                st.write(h.pregunta)
                                st.markdown("**Respuesta:**")
                                st.write(h.respuesta_final)
                                
                    # 🔗 MOSTRAR AUDIOS Y NOTAS DE NEUROVENTAS EN MACRO
                    st.markdown("---")
                    st.subheader("🎧 Notas y Audios Vinculados (CRM)")
                    
                    todos_los_prospectos = [prospect] + related
                    encontro_notas = False
                    for p in todos_los_prospectos:
                        if p.profile:
                            has_notes = bool(p.profile.notas_neuroventas)
                            has_audio = bool(p.profile.audio_path and os.path.exists(p.profile.audio_path))
                            if has_notes or has_audio:
                                encontro_notas = True
                                titulo = f"📌 {p.nombre} (RUT: {p.rut})" if p.id != prospect.id else f"📌 Este Cliente ({p.nombre})"
                                with st.expander(titulo):
                                    if p.id != prospect.id:
                                        st.caption("🔄 Extraído automáticamente de una entidad vinculada.")
                                    if has_notes:
                                        st.markdown(f"**Notas del Perfil:** {p.profile.notas_neuroventas}")
                                    if has_audio:
                                        st.markdown("**Audio Guardado:**")
                                        st.audio(p.profile.audio_path)
                                        
                    if not encontro_notas:
                        st.info("No se encontraron notas de audio ni análisis de neuroventas en el CRM para este cliente o sus relaciones.")
                else:
                    st.warning("Cliente no encontrado en la base de datos.")
            except Exception as e:
                st.error(f"Error cargando historial: {e}")
            finally:
                db.close()
        else:
            st.info("Selecciona un cliente para ver su historial de consultas macro.")

    with tab4:
        st.markdown("### 🔬 Inteligencia de Mercado & Playbooks")
        st.write("Combina las visiones macroeconómicas del mercado con estrategias de negocio e inteligencia sectorial.")
        from src.intelligence.market_analyst import render_market_intelligence
        render_market_intelligence()
